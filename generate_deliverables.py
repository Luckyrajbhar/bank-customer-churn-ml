from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

report_md = """# Bank Customer Churn Prediction Report

## 1. Project Overview
This AI-generated report describes a bank customer churn prediction solution that uses ensemble machine learning to identify customers at risk of leaving.

## 2. Problem Statement
Customer churn reduces revenue and increases acquisition cost. Predicting churn early enables banks to implement targeted retention strategies.

## 3. Data and Features
- Dataset includes credit score, age, tenure, balance, number of products, country, gender, and churn label.
- Removed irrelevant features such as customer ID.
- Encoded categorical variables: gender and country.
- Added `balance_per_product` as a new feature to capture customer engagement.

## 4. Preprocessing
- Standardized numerical features with `StandardScaler`.
- Applied one-hot encoding to categorical variables.
- Created a clean feature set for model training.

## 5. Model Architecture
- Trained an ensemble Voting Classifier.
- Combined SVM, Gradient Boosting, and XGBoost models.
- Used soft voting to aggregate prediction probabilities.

## 6. Application
- `app.py` is a Streamlit application for interactive churn prediction.
- Customers can be scored live based on input features.
- The app displays risk level and retention guidance.

## 7. Deliverables
- `Bank_Churn_Presentation.pptx`: AI-styled project presentation.
- `project_report.md`: Detailed written report.
- `generate_deliverables.py`: Script to regenerate deliverables anytime.

## 8. Recommendations
- Tune model hyperparameters and validate with cross-validation.
- Explore additional feature engineering and customer segmentation.
- Deploy the app for real-time retention analysis.
"""

with open("project_report.md", "w", encoding="utf-8") as f:
    f.write(report_md)

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Bank Customer Churn Prediction"
slide.placeholders[1].text = "AI-generated project presentation with a polished structure for business and technical audiences."

slides = [
    {
        "title": "Project Overview",
        "bullets": [
            "Predict whether bank customers will churn using machine learning.",
            "Use customer profile and financial behavior features.",
            "Built with Python, scikit-learn, XGBoost, and Streamlit."
        ]
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Customer churn causes lost revenue and higher acquisition costs.",
            "Objective: identify at-risk customers early.",
            "Enable retention actions with targeted interventions."
        ]
    },
    {
        "title": "Data & Feature Engineering",
        "bullets": [
            "Used demographic and account data such as age, balance, tenure.",
            "One-hot encoded gender and country features.",
            "Created `balance_per_product` to measure engagement."
        ]
    },
    {
        "title": "Modeling Approach",
        "bullets": [
            "Combined SVM, Gradient Boosting, and XGBoost in a Voting Classifier.",
            "Soft voting aggregates prediction probabilities for accuracy.",
            "Trained on scaled features for consistent model input."
        ]
    },
    {
        "title": "Interactive Application",
        "bullets": [
            "Streamlit app accepts customer details for real-time scoring.",
            "Displays churn probability and risk insights.",
            "Supports business decisions for customer retention."
        ]
    },
    {
        "title": "Results & Next Steps",
        "bullets": [
            "Deliver a deployable model and easy-to-use prediction app.",
            "Future work: hyperparameter tuning and richer features.",
            "Expand deployment options and integrate customer feedback."
        ]
    }
]

for slide_info in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_info["title"]
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(slide_info["bullets"]):
        if i == 0:
            p = body.paragraphs[0]
            p.text = bullet
        else:
            p = body.add_paragraph()
            p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

# Add a final slide with AI generation note
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "AI-Generated Deliverables"
body = slide.shapes.placeholders[1].text_frame
body.clear()
notes = [
    "This presentation was generated automatically using a Python script.",
    "It is designed to be clear, professional, and easy to present.",
    "Regenerate or customize the slides by editing generate_deliverables.py."
]
for i, note in enumerate(notes):
    if i == 0:
        p = body.paragraphs[0]
        p.text = note
    else:
        p = body.add_paragraph()
        p.text = note
    p.level = 0
    p.font.size = Pt(18)

prs.save("Bank_Churn_Presentation.pptx")
print("Updated Bank_Churn_Presentation.pptx and project_report.md")
